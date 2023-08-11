import copy
import logging
import os.path

import pptx
from pptx.util import Pt

from data_objects import Presentation, Topic, LAYOUT, Slide
from presentation_genrating_stage.presentation_generation import \
    GenerationHandler, Organizer
from templates import TEMPLATE
from utils import RESULTS_DIR, divide_sentence
from utils.Errors import NotFoundError
from utils.PresentationExportionUtils import replace_with_image, \
    CONTENT_PLACEHOLDER_IDX, DEFAULT_BULLET_POINT_FONT_SIZE, \
    NUM_OF_LINES_PER_SLIDE, estimate_line_count, get_font_size


class PresentationHandler:
    _presentation: Presentation
    _generator_handler: GenerationHandler
    _organizer: Organizer
    _BULLET_POINT_CHARACTERS_NUM_LIMIT: int
    _BULLET_POINTS_DELIMITERS: list[str]

    def __init__(self):
        self._presentation = Presentation()
        self._generator_handler = GenerationHandler()
        self._organizer = Organizer()
        self._BULLET_POINTS_DELIMITERS = ["while", "because", "hence"
                                          , "but", "where", "originally"
                                          , "including", "such", "primarily"]
        self._BULLET_POINT_CHARACTERS_NUM_LIMIT = 125

    @property
    def presentation(self):
        return self._presentation

    def reset(self):
        self._presentation = Presentation()

    def create_presentation(self, topic: Topic
                            , generators_names: list[str]
                            , params: dict[str, dict[str, object]]
                            , slides_number: int = -1) -> bool:
        # setting title
        try:
            self.presentation.title = topic.title

            self._generator_handler.generate_content(self.presentation, topic,
                                                     params, generators_names)

            logging.debug("keypoints generated successfully")
            logging.debug(self.presentation.all_keypoints)

            self._organizer.organize(self.presentation, topic, slides_number)
            logging.debug("presentation organized successfully")
            logging.debug(self.presentation.slides)

            return True
        except Exception as e:
            logging.error("Generation Failed")
            logging.exception(e)
            return False

    def export_presentation(self, path: str = RESULTS_DIR
                            , template: TEMPLATE = TEMPLATE.TEMPLATE_1):
        # fixme: this function is getting more complicated
        try:
            pr = pptx.Presentation(template.value)
            # adding title slide
            title_slide = pr.slides[0]
            title_slide.shapes.title.text = self.presentation.title

            # adding rest of slides
            for slide in self.presentation.slides:
                self.export_slide(slide, pr)

            # adding conclusion slide
            conclusion_slide_layout = pr.slide_layouts[LAYOUT.TITLE.value]
            conclusion_slide = pr.slides.add_slide(conclusion_slide_layout)
            conclusion_slide.shapes.title.text = "Conclusion"

            save_path = os.path.join(path, self.presentation.title + ".pptx")
            pr.save(save_path)

            logging.debug("Presentation exported successfully")
        except Exception as e:
            logging.error("Exporting FAILED")
            logging.exception(e)

    def get_presentation(self):
        return copy.deepcopy(self.presentation)

    def export_slide(self, slide: Slide, presentation_to_export):
        bullet_points_text = [str(k.data) for k in slide.keypoints]
        # handling first slide
        slide_layout = presentation_to_export.slide_layouts[slide.layout.value]
        m_slide = presentation_to_export.slides.add_slide(slide_layout)
        m_slide.shapes.title.text = slide.title
        # choosing the right placeholder for bullet points
        if slide.layout == LAYOUT.TITLE_AND_CONTENT:
            holder_index = 1
        else:
            holder_index = 2
        bullet_point_box = m_slide.shapes.placeholders[holder_index]
        # adding images if existed to each exported slide
        if slide.layout == LAYOUT.PICTURE_CONTENT_WITH_CAPTION:
            image_box = m_slide.placeholders[CONTENT_PLACEHOLDER_IDX]
            replace_with_image(str(slide.content.data)
                               , image_box, m_slide)
        if len(bullet_points_text) > 0:
            self._insert_bullet_point(bullet_point_box
                                      , bullet_points_text[0]
                                      , slide.keypoints[0].level)
        # handling rest of slides
        for idx, bullet_point in enumerate(bullet_points_text[1:]):
            # if there is overflow or the deck is empty
            if estimate_line_count(bullet_point_box) > NUM_OF_LINES_PER_SLIDE:
                slide_layout = presentation_to_export.slide_layouts[
                    slide.layout.value]
                m_slide = presentation_to_export.slides.add_slide(slide_layout)
                m_slide.shapes.title.text = slide.title
                # choosing the right placeholder for bullet points
                if slide.layout == LAYOUT.TITLE_AND_CONTENT:
                    holder_index = 1
                else:
                    holder_index = 2
                bullet_point_box = m_slide.shapes.placeholders[holder_index]
                # adding images if existed to each exported slide
                if slide.layout == LAYOUT.PICTURE_CONTENT_WITH_CAPTION:
                    image_box = m_slide.placeholders[CONTENT_PLACEHOLDER_IDX]
                    replace_with_image(str(slide.content.data)
                                       , image_box, m_slide)
            if bullet_point_box is not None:
                self._insert_bullet_point(bullet_point_box
                                          , bullet_point
                                          , slide.keypoints[idx].level)

    def _get_paragraph(self, box):
        """
        in python-pptx, paragraphs is not initially empty, it has one
        empty paragraph and this function to handle this case :param
        self: :param box: a shape that has a text_frame :return: adds new
        paragraph and returned it or use the already existing one
        """
        if len(box.text_frame.paragraphs[0].text) == 0:
            return box.text_frame.paragraphs[0]
        else:
            return box.text_frame.add_paragraph()

    def _insert_bullet_point(self, box, bullet_point: str, level: int = -1):
        # divide bullet point.
        bullets = []
        if len(bullet_point) > self._BULLET_POINT_CHARACTERS_NUM_LIMIT:
            bullets = divide_sentence(bullet_point
                                      , self._BULLET_POINTS_DELIMITERS)
        else:
            bullets.append(bullet_point)
        # add bullet point with first level 1 and rest level 2
        for idx, bullet in enumerate(bullets):
            paragraph = self._get_paragraph(box)
            paragraph.text = bullet
            if level == -1:
                if idx == 0:
                    paragraph.level = 0
                else:
                    paragraph.level = 1
            else:
                paragraph.level = level

        # adjusting the font which is important for counting the font later.
        for p in box.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(get_font_size(p.level))

